"""Build the leadership deck.

Light ground, because a projector and a printer both handle it and a dark
deck gambles on the room. Arial throughout: IBM Plex is not installed on
every machine this will be opened on, and a deck whose design depends on a
font the audience lacks is a broken deck. The character comes from spacing,
hierarchy and a palette that means something — teal is ours and fixed,
amber is the client's and varies, red is a refusal.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK        = RGBColor(0x12, 0x1A, 0x19)
INK2       = RGBColor(0x3A, 0x48, 0x46)
MUTED      = RGBColor(0x61, 0x70, 0x6D)
LINE       = RGBColor(0xD2, 0xDA, 0xD8)
GROUND     = RGBColor(0xF1, 0xF4, 0xF3)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEAL       = RGBColor(0x0E, 0x6F, 0x62)
TEAL_SOFT  = RGBColor(0xDD, 0xEB, 0xE8)
AMBER      = RGBColor(0x9C, 0x64, 0x13)
AMBER_SOFT = RGBColor(0xF4, 0xE9, 0xD6)
CRIT       = RGBColor(0x9B, 0x35, 0x27)
CRIT_SOFT  = RGBColor(0xF5, 0xE1, 0xDD)

FONT = "Arial"
W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)                      # page margin
CONTENT_W = W - 2 * M


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs, ground=GROUND):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = ground
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    """runs: list of (text, size, bold, colour, space_after_pt)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (body, size, bold, colour, after) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(after)
        run = para.add_run(); run.text = body
        f = run.font
        f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, colour
    return box


def eyebrow(slide, label):
    text(slide, M, Inches(0.55), CONTENT_W, Inches(0.3),
         [(label.upper(), 11, True, MUTED, 0)])


def title(slide, heading, sub=None):
    """Box sized to what the text needs.

    A nominal 1.5" box for a one-line heading does not collide when rendered
    but does collide as geometry — and would collide for real the moment a
    heading wrapped to two lines. Sizing it honestly means the overlap check
    is telling the truth.
    """
    runs = [(heading, 34, True, INK, 6 if sub else 0)]
    if sub:
        runs.append((sub, 16, False, INK2, 0))
    height = Inches(1.45) if sub else Inches(0.95)
    text(slide, M, Inches(1.0), CONTENT_W, height, runs, spacing=1.05)


def rule(slide, y, colour=LINE, thickness=1.0, x=M, w=None):
    w = w or CONTENT_W
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(thickness))
    ln.fill.solid(); ln.fill.fore_color.rgb = colour
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def card(slide, x, y, w, h, *, fill=WHITE, border=LINE, width=1.0, dash=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(width)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    box.shadow.inherit = False
    box.text_frame.text = ""
    return box


def table(slide, x, y, w, rows, widths, *, header=True, row_h=Inches(0.52), size=13):
    """rows[0] is the header when header=True."""
    n, cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n, cols, x, y, w, row_h * n)
    tbl = shape.table
    total = sum(widths)
    for c, frac in enumerate(widths):
        tbl.columns[c].width = Emu(int(w * frac / total))
    for r, row in enumerate(rows):
        tbl.rows[r].height = row_h
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.14)
            cell.margin_top = cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = GROUND if (header and r == 0) else WHITE
            para = cell.text_frame.paragraphs[0]
            para.line_spacing = 1.1
            run = para.add_run(); run.text = str(val)
            f = run.font
            f.name = FONT
            f.size = Pt(11 if (header and r == 0) else size)
            f.bold = header and r == 0
            f.color.rgb = MUTED if (header and r == 0) else INK
    return tbl


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body
