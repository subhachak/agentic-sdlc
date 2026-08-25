"""Static checks on the built deck.

No renderer is available here, so these substitute for looking at it: nothing
outside the slide, no two text boxes on top of each other, and no text that
cannot fit the box it was put in. Text height is estimated from Arial's
average advance width — rough, but it catches the gross overflow that a deck
built by arithmetic actually suffers from.
"""
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400
prs = Presentation("/Users/subh/MyApps/agentic-sdlc/docs/agentic-sdlc-leadership.pptx")
SW, SH = prs.slide_width, prs.slide_height

def est_height(tf, width_emu):
    """Rough rendered height in EMU."""
    width_pt = width_emu / EMU_IN * 72
    total = 0.0
    for p in tf.paragraphs:
        runs = p.runs
        if not runs:
            continue
        size = max((r.font.size.pt if r.font.size else 18) for r in runs)
        chars = sum(len(r.text) for r in runs)
        # Arial average advance ≈ 0.52em; hard newlines force a break.
        hard = sum(r.text.count("\n") for r in runs)
        per_line = max(1, int(width_pt / (size * 0.52)))
        lines = max(1, -(-chars // per_line)) + hard
        spacing = p.line_spacing if isinstance(p.line_spacing, float) else 1.15
        total += lines * size * spacing * 1.02
        total += (p.space_after.pt if p.space_after else 0)
    return total / 72 * EMU_IN

problems = []
for i, slide in enumerate(prs.slides, 1):
    boxes = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        r, b = sh.left + sh.width, sh.top + sh.height
        if sh.left < -1000 or sh.top < -1000 or r > SW + 1000 or b > SH + 1000:
            problems.append(f"s{i}: '{(sh.name or '')[:18]}' outside slide "
                            f"({sh.left/EMU_IN:.2f},{sh.top/EMU_IN:.2f} → {r/EMU_IN:.2f},{b/EMU_IN:.2f})")
        if sh.has_text_frame and sh.text_frame.text.strip():
            need = est_height(sh.text_frame, sh.width)
            if need > sh.height * 1.35:
                problems.append(f"s{i}: text needs {need/EMU_IN:.2f}\" in {sh.height/EMU_IN:.2f}\" box "
                                f"— \"{sh.text_frame.text.strip()[:44]}\"")
            boxes.append((sh.left, sh.top, sh.width, sh.height, sh.text_frame.text.strip()[:26]))

    for a in range(len(boxes)):
        for c in range(a + 1, len(boxes)):
            x1, y1, w1, h1, t1 = boxes[a]; x2, y2, w2, h2, t2 = boxes[c]
            ox = min(x1 + w1, x2 + w2) - max(x1, x2)
            oy = min(y1 + h1, y2 + h2) - max(y1, y2)
            if ox > EMU_IN * 0.05 and oy > EMU_IN * 0.05:
                problems.append(f"s{i}: text overlap \"{t1}\" ∩ \"{t2}\"")

    if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
        problems.append(f"s{i}: no speaker notes")

print(f"{len(prs.slides)} slides checked")
if problems:
    print(f"\n{len(problems)} issue(s):")
    for p in problems[:24]:
        print("  -", p)
else:
    print("no issues found")
